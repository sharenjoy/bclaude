#!/usr/bin/env python3
"""
PPT Maker - PPTX 打包腳本
=========================
將圖片目錄打包成 PPTX 檔案。

使用方式：
    python build_pptx.py --images ./images --output output.pptx
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("錯誤：請安裝 python-pptx")
    print("執行：pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print("錯誤：請安裝 Pillow")
    print("執行：pip install Pillow")
    sys.exit(1)


def build_pptx(images_dir, output_file):
    """將圖片打包成 PPTX"""
    images_path = Path(images_dir)
    output_path = Path(output_file)

    if not images_path.exists():
        print(f"錯誤：圖片目錄不存在 - {images_dir}")
        return False

    # 收集圖片
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp'}
    images = sorted([
        f for f in images_path.iterdir()
        if f.suffix.lower() in image_extensions
    ])

    if not images:
        print(f"錯誤：找不到圖片 - {images_dir}")
        return False

    print(f"找到 {len(images)} 張圖片")

    # 建立 PPTX
    prs = Presentation()

    # 設定 16:9 尺寸
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches

    # 取得空白版面
    blank_layout = prs.slide_layouts[6]

    for i, image_path in enumerate(images, 1):
        print(f"處理 {i}/{len(images)}：{image_path.name}")

        # 新增空白投影片
        slide = prs.slides.add_slide(blank_layout)

        # 取得圖片尺寸
        with Image.open(image_path) as img:
            img_width, img_height = img.size

        # 計算縮放比例，保持比例
        img_ratio = img_width / img_height
        slide_ratio = slide_width / slide_height

        if img_ratio > slide_ratio:
            width = Inches(slide_width)
            height = Inches(slide_width / img_ratio)
            left = Inches(0)
            top = Inches((slide_height - slide_width / img_ratio) / 2)
        else:
            height = Inches(slide_height)
            width = Inches(slide_height * img_ratio)
            left = Inches((slide_width - slide_height * img_ratio) / 2)
            top = Inches(0)

        # 加入圖片
        slide.shapes.add_picture(str(image_path), left, top, width, height)

    # 確保輸出目錄存在
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # 儲存
    prs.save(str(output_path))
    print(f"PPTX 已儲存：{output_path}")
    print(f"共 {len(images)} 頁")

    return True


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPT Maker - PPTX 打包")
    parser.add_argument("--images", required=True, help="圖片目錄")
    parser.add_argument("--output", required=True, help="輸出 PPTX 路徑")

    args = parser.parse_args()
    success = build_pptx(args.images, args.output)
    sys.exit(0 if success else 1)
